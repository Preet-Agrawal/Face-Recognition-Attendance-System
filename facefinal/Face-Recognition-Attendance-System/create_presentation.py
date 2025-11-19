"""
Script to create a PowerPoint presentation for Face Recognition Attendance System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title = title_frame.add_paragraph()
    title.text = "Face Recognition Attendance System"
    title.font.size = Pt(44)
    title.font.bold = True
    title.font.color.rgb = RGBColor(0, 51, 102)
    title.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle = subtitle_frame.add_paragraph()
    subtitle.text = "AI-Powered Automated Attendance Management"
    subtitle.font.size = Pt(24)
    subtitle.font.color.rgb = RGBColor(102, 102, 102)
    subtitle.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    
    content = slide.placeholders[1].text_frame
    content.text = "Traditional attendance systems face multiple challenges:"
    
    points = [
        "Time-consuming manual roll calls",
        "Proxy attendance and fraud",
        "Human errors in record keeping",
        "Difficulty in tracking attendance patterns",
        "Paper-based systems are inefficient"
    ]
    
    for point in points:
        p = content.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 3: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Our Solution"
    
    content = slide.placeholders[1].text_frame
    content.text = "An intelligent face recognition system that:"
    
    points = [
        "Automatically marks attendance using facial recognition",
        "Provides real-time attendance tracking",
        "Prevents proxy attendance with 99.6% accuracy",
        "Offers web-based interface for easy access",
        "Exports attendance reports to Excel",
        "Stores data securely in a database"
    ]
    
    for point in points:
        p = content.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"
    
    content = slide.placeholders[1].text_frame
    
    features = [
        ("Real-time Face Recognition", "MTCNN + FaceNet for 99.6% accuracy"),
        ("Multi-angle Registration", "3-7 images per person for robust recognition"),
        ("Quality Assessment", "Ensures high-quality face captures"),
        ("Anti-spoofing Detection", "Prevents photo/video-based fraud"),
        ("Web Dashboard", "Modern Flask-based interface"),
        ("Excel Export", "Easy attendance report generation")
    ]
    
    for feature, desc in features:
        p = content.add_paragraph()
        p.text = f"{feature}: {desc}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 5: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1].text_frame
    
    tech_stack = [
        ("Backend", "Python, Flask"),
        ("Face Detection", "MTCNN (Multi-task Cascaded CNN)"),
        ("Face Recognition", "FaceNet (InceptionResnetV1)"),
        ("Deep Learning", "PyTorch, facenet-pytorch"),
        ("Computer Vision", "OpenCV, dlib"),
        ("Database", "SQLite with pandas"),
        ("Frontend", "HTML, CSS, JavaScript"),
        ("Export", "openpyxl for Excel generation")
    ]
    
    for tech, tools in tech_stack:
        p = content.add_paragraph()
        p.text = f"{tech}: {tools}"
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 6: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    
    content = slide.placeholders[1].text_frame
    
    components = [
        "Web Interface (Flask)",
        "├── Camera Module (OpenCV)",
        "├── Face Detection (MTCNN)",
        "├── Face Recognition (FaceNet)",
        "├── Database Manager (SQLite)",
        "└── Export Module (Excel)",
        "",
        "Data Flow:",
        "Camera → Detection → Recognition → Database → Reports"
    ]
    
    for component in components:
        p = content.add_paragraph()
        p.text = component
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = 'Courier New'
    
    # Slide 7: Core Modules
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Core Modules"
    
    content = slide.placeholders[1].text_frame
    
    modules = [
        ("app.py", "Flask web server with REST API endpoints"),
        ("production_face_recognition.py", "Face detection & recognition engine"),
        ("database_manager.py", "SQLite database operations"),
        ("templates/", "Web interface HTML templates")
    ]
    
    for module, desc in modules:
        p = content.add_paragraph()
        p.text = f"{module}"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        
        p2 = content.add_paragraph()
        p2.text = desc
        p2.level = 1
        p2.font.size = Pt(14)
    
    # Slide 8: How It Works
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How It Works"
    
    content = slide.placeholders[1].text_frame
    content.text = "Registration Process:"
    
    reg_steps = [
        "Student enters details (ID, name, email, etc.)",
        "System captures 3-7 face images from different angles",
        "Quality assessment ensures clear images",
        "FaceNet generates 512D embeddings",
        "Embeddings stored in database"
    ]
    
    for step in reg_steps:
        p = content.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(16)
    
    p = content.add_paragraph()
    p.text = "\nAttendance Marking:"
    p.level = 0
    p.font.size = Pt(18)
    
    att_steps = [
        "Camera captures live video feed",
        "MTCNN detects faces in real-time",
        "FaceNet extracts face embeddings",
        "Compares with database using cosine similarity",
        "Auto-marks attendance if match found (>65% similarity)"
    ]
    
    for step in att_steps:
        p = content.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 9: API Endpoints
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "REST API Endpoints"
    
    content = slide.placeholders[1].text_frame
    
    endpoints = [
        "/api/start_camera - Start camera feed",
        "/api/register/start - Begin registration",
        "/api/register/capture - Capture face image",
        "/api/register/complete - Finish registration",
        "/api/students - Get all students",
        "/api/attendance/today - Today's attendance",
        "/api/export/today - Export to Excel",
        "/api/stats - System statistics"
    ]
    
    for endpoint in endpoints:
        p = content.add_paragraph()
        p.text = endpoint
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = 'Courier New'
    
    # Slide 10: Performance Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Metrics"
    
    content = slide.placeholders[1].text_frame
    
    metrics = [
        ("Recognition Accuracy", "99.6% (FaceNet on VGGFace2)"),
        ("Recognition Threshold", "65% cosine similarity"),
        ("Quality Threshold", "40% minimum quality score"),
        ("Processing Speed", "~30 FPS real-time"),
        ("Min Face Size", "60x60 pixels"),
        ("Embedding Dimension", "512D vector space"),
        ("Anti-spoofing", "Enabled by default")
    ]
    
    for metric, value in metrics:
        p = content.add_paragraph()
        p.text = f"{metric}: {value}"
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 11: Security Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security Features"
    
    content = slide.placeholders[1].text_frame
    
    security = [
        "Anti-spoofing detection prevents photo/video attacks",
        "Quality-based filtering ensures genuine captures",
        "Multi-angle registration for robust verification",
        "Ensemble verification with multiple embeddings",
        "Secure database storage with pickle serialization",
        "No face images stored, only embeddings",
        "CORS enabled for secure web access"
    ]
    
    for feature in security:
        p = content.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 12: Advantages
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Advantages"
    
    content = slide.placeholders[1].text_frame
    
    advantages = [
        "Contactless and hygienic",
        "Saves time - no manual roll calls",
        "Eliminates proxy attendance",
        "Real-time attendance tracking",
        "Automated report generation",
        "Scalable to large institutions",
        "Easy to use web interface",
        "Cost-effective solution"
    ]
    
    for adv in advantages:
        p = content.add_paragraph()
        p.text = adv
        p.level = 0
        p.font.size = Pt(20)
    
    # Slide 13: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Use Cases"
    
    content = slide.placeholders[1].text_frame
    
    use_cases = [
        ("Educational Institutions", "Schools, colleges, universities"),
        ("Corporate Offices", "Employee attendance tracking"),
        ("Training Centers", "Workshop and seminar attendance"),
        ("Events", "Conference and event check-ins"),
        ("Laboratories", "Research lab access control"),
        ("Libraries", "Visitor tracking and management")
    ]
    
    for use_case, desc in use_cases:
        p = content.add_paragraph()
        p.text = f"{use_case}: {desc}"
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 14: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    
    content = slide.placeholders[1].text_frame
    
    enhancements = [
        "Mobile app integration",
        "SMS/Email notifications to parents/managers",
        "Advanced analytics and insights",
        "Integration with existing ERP systems",
        "Multi-camera support",
        "Cloud deployment for remote access",
        "Facial expression analysis for engagement tracking",
        "GPU acceleration for faster processing"
    ]
    
    for enhancement in enhancements:
        p = content.add_paragraph()
        p.text = enhancement
        p.level = 0
        p.font.size = Pt(18)
    
    # Slide 15: Installation & Setup
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Installation & Setup"
    
    content = slide.placeholders[1].text_frame
    
    steps = [
        "1. Clone the repository",
        "2. Install dependencies:",
        "   pip install -r requirements.txt",
        "3. Run the application:",
        "   python app.py",
        "4. Open browser:",
        "   http://localhost:8000",
        "5. Register students and start marking attendance!"
    ]
    
    for step in steps:
        p = content.add_paragraph()
        p.text = step
        p.level = 0
        p.font.size = Pt(18)
        if step.startswith("   "):
            p.font.name = 'Courier New'
    
    # Slide 16: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1].text_frame
    
    conclusion = [
        "✓ Automated, accurate, and efficient attendance system",
        "✓ State-of-the-art AI technology (99.6% accuracy)",
        "✓ User-friendly web interface",
        "✓ Secure and scalable solution",
        "✓ Eliminates manual errors and fraud",
        "✓ Ready for deployment in real-world scenarios"
    ]
    
    for point in conclusion:
        p = content.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
    
    # Slide 17: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    thank_you_frame = thank_you_box.text_frame
    thank_you = thank_you_frame.add_paragraph()
    thank_you.text = "Thank You!"
    thank_you.font.size = Pt(54)
    thank_you.font.bold = True
    thank_you.font.color.rgb = RGBColor(0, 51, 102)
    thank_you.alignment = PP_ALIGN.CENTER
    
    questions_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.8))
    questions_frame = questions_box.text_frame
    questions = questions_frame.add_paragraph()
    questions.text = "Questions?"
    questions.font.size = Pt(28)
    questions.font.color.rgb = RGBColor(102, 102, 102)
    questions.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Face_Recognition_Attendance_System_Presentation.pptx')
    print("✅ Presentation created successfully!")
    print("📄 File: Face_Recognition_Attendance_System_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
